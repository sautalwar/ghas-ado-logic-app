"""Generate GHAS-ADO Logic App Demo PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color constants
DARK_BG = RGBColor(0x1B, 0x1F, 0x23)
ACCENT_BLUE = RGBColor(0x00, 0x78, 0xD4)
ACCENT_GREEN = RGBColor(0x2E, 0xA0, 0x43)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
MID_GRAY = RGBColor(0xA0, 0xA0, 0xA0)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)

screenshots_dir = "docs/screenshots"


def add_bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_bullets(slide, left, top, width, height, items,
                font_size=16, color=WHITE):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return txBox


def add_img(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        kwargs = {}
        if width:
            kwargs["width"] = Inches(width)
        if height:
            kwargs["height"] = Inches(height)
        slide.shapes.add_picture(path, Inches(left), Inches(top), **kwargs)


# ===================== SLIDE 1: Title =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 1, 1.5, 11, 1.2,
             "GHAS \u2192 Azure DevOps", 44, ACCENT_BLUE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 2.8, 11, 1,
             "Automated Vulnerability Tracking", 36, WHITE, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.2, 11, 0.8,
             "Seamless integration between GitHub Advanced Security "
             "and Azure DevOps Work Items",
             18, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.5, 11, 0.5,
             "Learfield Customer Demo", 20, ORANGE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 6.5, 11, 0.5,
             "Powered by Azure Logic Apps  |  Bicep IaC  |  GitHub Webhooks",
             14, MID_GRAY, False, PP_ALIGN.CENTER)

# ===================== SLIDE 2: The Problem =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "The Problem", 36, ACCENT_BLUE, True)
add_bullets(slide, 0.8, 1.5, 5.5, 5.5, [
    "\u26a0\ufe0f  Security vulnerabilities detected in GHAS go untracked",
    "\U0001f504  Manual process: Open GHAS \u2192 Go to ADO \u2192 Create work item \u2192 Link back",
    "\u23f1\ufe0f  Time-consuming and error-prone \u2014 developers skip it",
    "\U0001f517  No traceability between security alerts and dev work",
    "\u274c  Work items stay open even after vulnerabilities are fixed",
    "\U0001f4ca  Security team lacks visibility into remediation progress",
], 18, LIGHT_GRAY)
add_text_box(slide, 7, 1.5, 5.5, 3,
             '"We need this to happen automatically.\n'
             'When GHAS finds something,\n'
             'ADO should know about it instantly.\n'
             'No human in the loop."',
             22, ORANGE, True)
add_text_box(slide, 7, 4.5, 5.5, 1,
             "\u2014 Customer Feedback (Learfield)", 16, MID_GRAY, True)

# ===================== SLIDE 3: The Solution =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 11, 0.8,
             "The Solution: Azure Logic App", 36, ACCENT_BLUE, True)
add_bullets(slide, 0.8, 1.5, 11, 2.5, [
    "\U0001f514  GitHub webhook triggers on every GHAS alert",
    "\u26a1  Azure Logic App processes the event in real-time",
    "\U0001f4cb  ADO work item created automatically with full metadata",
    "\u2705  Auto-closes work item when vulnerability is resolved",
    "\U0001f3f7\ufe0f  Deduplication via tags \u2014 no duplicate work items",
], 17, LIGHT_GRAY)
add_text_box(slide, 0.8, 4.2, 11, 1.2,
             "GitHub GHAS  \u2500\u2500webhook\u2500\u2500\u25b6  "
             "Azure Logic App  \u2500\u2500REST API\u2500\u2500\u25b6  "
             "Azure DevOps\n"
             "(Code Scanning, Dependabot, Secret Scanning)"
             "       (Consumption tier)"
             "              (Work Items + Tags)",
             14, ACCENT_GREEN, True, PP_ALIGN.CENTER)
add_bullets(slide, 0.8, 5.8, 11, 1.5, [
    "Supports all 3 GHAS alert types  \u2022  "
    "Infrastructure as Code (Bicep)  \u2022  Zero ongoing maintenance",
], 15, MID_GRAY)

# ===================== SLIDE 4: Architecture =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 11, 0.8,
             "Architecture Overview", 36, ACCENT_BLUE, True)

arch_text = (
    "\u250c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500"
    "\u2500\u2500\u2500\u2500\u2510        "
    "\u250c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500"
    "\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500"
    "\u2500\u2500\u2500\u2510        "
    "\u250c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500"
    "\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2510\n"
    "\u2502   GitHub Repo    \u2502        "
    "\u2502     Azure Logic App          \u2502        "
    "\u2502   Azure DevOps       \u2502\n"
    "\u2502                  \u2502        "
    "\u2502                              \u2502        "
    "\u2502                      \u2502\n"
    "\u2502  GHAS Alerts:    \u2502\u2500\u2500\u2500\u2500\u2500\u25b6"
    "\u2502  1. Receive Webhook          \u2502\u2500\u2500\u2500\u2500\u2500\u25b6"
    "\u2502  Query existing WIs  \u2502\n"
    "\u2502  \u2022 Code Scanning \u2502 POST  "
    "\u2502  2. Parse alert metadata     \u2502 WIQL  "
    "\u2502  (dedup check)       \u2502\n"
    "\u2502  \u2022 Dependabot    \u2502       "
    "\u2502  3. Normalize fields         \u2502       "
    "\u2502                      \u2502\n"
    "\u2502  \u2022 Secret Scan   \u2502       "
    "\u2502  4. Check: create or close?  \u2502       "
    "\u2502  Create new Issue    \u2502\n"
    "\u2502                  \u2502       "
    "\u2502  5. Dedup via GHAS tag       \u2502 PATCH "
    "\u2502  OR close existing   \u2502\n"
    "\u2502  Webhook fires   \u2502       "
    "\u2502  6. Create / Close WI        \u2502       "
    "\u2502  with comment        \u2502\n"
    "\u2514\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500"
    "\u2500\u2500\u2500\u2500\u2518        "
    "\u2514\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500"
    "\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500"
    "\u2500\u2500\u2500\u2518        "
    "\u2514\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500"
    "\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518"
)
add_text_box(slide, 0.5, 1.4, 12, 4, arch_text, 12, ACCENT_GREEN, True, PP_ALIGN.CENTER)

add_bullets(slide, 0.8, 5.2, 11, 2, [
    "Work Item Fields: Title, Description (HTML), Severity tag, "
    "Repo, Branch, File:Line, Hyperlink to GHAS alert",
    "Tag format: GHAS-{owner}-{repo}-{alertNumber} \u2014 "
    "enables deduplication and auto-close lookup",
], 14, LIGHT_GRAY)

# ===================== SLIDE 5: Azure Portal =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7,
             "Step 1: Create Logic App in Azure Portal", 32, ACCENT_BLUE, True)
add_img(slide, f"{screenshots_dir}/01-azure-portal-home.png", 0.3, 1.1, 6.2)
add_img(slide, f"{screenshots_dir}/02-create-resource-logic-app.png", 6.8, 1.1, 6.2)
add_text_box(slide, 0.3, 6.3, 6, 0.5,
             "Azure Portal Home", 12, MID_GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 6.8, 6.3, 6, 0.5,
             "Search & Select Logic App", 12, MID_GRAY, False, PP_ALIGN.CENTER)

# ===================== SLIDE 6: Logic App Config =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7,
             "Step 2: Configure Logic App (Consumption)", 32, ACCENT_BLUE, True)
add_img(slide, f"{screenshots_dir}/03-select-consumption-plan.png", 0.3, 1.1, 6.2)
add_img(slide, f"{screenshots_dir}/04-logic-app-basics-form.png", 6.8, 1.1, 6.2)
add_text_box(slide, 0.3, 6.3, 6, 0.5,
             "Select Consumption Plan", 12, MID_GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 6.8, 6.3, 6, 0.5,
             "Configure Basics", 12, MID_GRAY, False, PP_ALIGN.CENTER)

# ===================== SLIDE 7: Logic App Overview =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7,
             "Step 3: Logic App Deployed", 32, ACCENT_BLUE, True)
add_img(slide, f"{screenshots_dir}/05-logic-app-overview.png", 1.5, 1.1, 10)
add_text_box(slide, 1.5, 6.5, 10, 0.5,
             "Logic App Overview \u2014 Status: Enabled, 1 Trigger, 22 Actions",
             14, MID_GRAY, False, PP_ALIGN.CENTER)

# ===================== SLIDE 8: Designer Full Workflow =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7,
             "Workflow Designer: Full Flow", 32, ACCENT_BLUE, True)
add_img(slide, f"{screenshots_dir}/06-designer-full-workflow.png", 1.5, 1.1, 10)
add_text_box(slide, 1.5, 6.5, 10, 0.5,
             "14 Compose actions normalize metadata \u2192 "
             "Condition routes to Create or Close path",
             14, MID_GRAY, False, PP_ALIGN.CENTER)

# ===================== SLIDE 9: Condition & Branches =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7,
             "Workflow: Create vs Close Branches", 32, ACCENT_BLUE, True)
add_img(slide, f"{screenshots_dir}/07-designer-condition-expanded.png",
        0.3, 1.1, 6.2)
add_img(slide, f"{screenshots_dir}/09-designer-both-branches.png",
        6.8, 1.1, 6.2)
add_text_box(slide, 0.3, 6.3, 6, 0.5,
             'Condition: Is action "created"?',
             12, MID_GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 6.8, 6.3, 6, 0.5,
             "True: Create WI  |  False: Close WI",
             12, MID_GRAY, False, PP_ALIGN.CENTER)

# ===================== SLIDE 10: Trigger & HTTP =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7,
             "Webhook Trigger & ADO Integration", 32, ACCENT_BLUE, True)
add_img(slide, f"{screenshots_dir}/10-trigger-http-request.png", 0.3, 1.1, 6.2)
add_img(slide, f"{screenshots_dir}/11-http-query-existing-workitem.png",
        6.8, 1.1, 6.2)
add_text_box(slide, 0.3, 6.3, 6, 0.5,
             "HTTP Trigger (Webhook URL)", 12, MID_GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 6.8, 6.3, 6, 0.5,
             "WIQL Query for Deduplication", 12, MID_GRAY, False, PP_ALIGN.CENTER)

# ===================== SLIDE 11: ADO Work Items =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7,
             "Live Demo: Work Items Auto-Created in ADO",
             32, ACCENT_GREEN, True)
add_img(slide, f"{screenshots_dir}/13-ado-work-items-list.png", 1.5, 1.1, 10)
add_text_box(slide, 1.5, 6.5, 10, 0.5,
             "ADO Board showing GHAS-created issues (#10, #11) "
             "alongside seeded Learfield demo data",
             14, MID_GRAY, False, PP_ALIGN.CENTER)

# ===================== SLIDE 12: E2E Results =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 11, 0.8,
             "End-to-End Test Results", 36, ACCENT_GREEN, True)
add_bullets(slide, 0.8, 1.5, 11, 5.5, [
    "\u2705  Code Scanning Alert \u2192 Created WI #9: "
    "[GHAS-CodeScan] SQL Injection vulnerability",
    "      Tags: GHAS, CodeScanning, high, GHAS-learfield-corp-fan-portal-42",
    "",
    "\u2705  Dependabot Alert \u2192 Created WI #10: "
    "[GHAS-Dependabot] Prototype Pollution in lodash",
    "      Tags: GHAS, Dependabot, critical, GHAS-learfield-corp-fan-portal-7",
    "",
    "\u2705  Secret Scanning Alert \u2192 Created WI #11: "
    "[GHAS-Secret] Azure Storage Account Key",
    "      Tags: GHAS, SecretScanning, critical, GHAS-learfield-corp-fan-portal-3",
    "",
    "\u2705  Auto-Close: Sent \"fixed\" webhook \u2192 "
    "WI #9 transitioned to \"Done\" automatically",
    '      History comment: "Auto-closed by GHAS-ADO Sync: '
    'Vulnerability marked as fixed"',
    "",
    "\u2705  Deduplication: Duplicate webhooks do NOT create "
    "duplicate work items",
], 16, LIGHT_GRAY)

# ===================== SLIDE 13: Metadata =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 11, 0.8,
             "Work Item Metadata \u2014 Auto-Populated", 36, ACCENT_BLUE, True)

fields = [
    ("Field", "Source", "Example"),
    ("Title", "Alert type + rule description",
     "[GHAS-CodeScan] SQL Injection vulnerability"),
    ("Description", "HTML with full details",
     "Severity, repo, branch, file, line, alert link"),
    ("Tags", "GHAS + type + severity + dedup tag",
     "GHAS, CodeScanning, high, GHAS-org-repo-42"),
    ("Hyperlink", "Alert URL from webhook",
     "https://github.com/org/repo/security/..."),
    ("State", "Created \u2192 To Do / Fixed \u2192 Done",
     "Automatic lifecycle management"),
]

y = 1.5
for i, (field, source, example) in enumerate(fields):
    color = ORANGE if i == 0 else LIGHT_GRAY
    bold = i == 0
    sz = 15 if i == 0 else 14
    add_text_box(slide, 0.8, y, 2.5, 0.45, field, sz, color, bold)
    add_text_box(slide, 3.5, y, 4, 0.45, source, sz, color, bold)
    add_text_box(slide, 7.8, y, 5, 0.45, example, sz, color, bold)
    y += 0.55

add_bullets(slide, 0.8, 5.0, 11, 2, [
    "All 3 GHAS alert types supported:",
    "\u2022 Code Scanning \u2014 rule ID, severity, file path, line number",
    "\u2022 Dependabot \u2014 package name, advisory summary, manifest path",
    "\u2022 Secret Scanning \u2014 secret type, always tagged as critical",
], 15, MID_GRAY)

# ===================== SLIDE 14: IaC =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 11, 0.8,
             "Infrastructure as Code (Bicep)", 36, ACCENT_BLUE, True)
add_img(slide, f"{screenshots_dir}/12-logic-app-code-view.png", 0.3, 1.3, 7)
add_bullets(slide, 7.8, 1.3, 5, 5.5, [
    "\U0001f4c1 Repository Structure:",
    "   infra/main.bicep",
    "   infra/modules/logic-app.bicep",
    "   infra/workflows/ghas-to-ado.json",
    "   infra/parameters.json",
    "   scripts/deploy.ps1",
    "   scripts/setup-webhooks.ps1",
    "",
    "\U0001f680 One-command deploy:",
    "   az deployment group create \\",
    "     --template-file infra/main.bicep",
    "",
    "\U0001f504 Repeatable & version-controlled",
    "\U0001f512 PATs stored as parameters, not hardcoded",
], 13, LIGHT_GRAY)

# ===================== SLIDE 15: Value =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 11, 0.8,
             "Value Delivered", 36, ACCENT_GREEN, True)

values = [
    ("\u23f1\ufe0f  Time Saved",
     "Eliminates manual work item creation \u2014 "
     "saves 5-10 min per vulnerability"),
    ("\U0001f517  Full Traceability",
     "Every GHAS alert linked to an ADO work item with hyperlink back"),
    ("\U0001f916  Zero Human Effort",
     "No developer action needed \u2014 Logic App handles everything"),
    ("\u2705  Auto-Close",
     "Work items automatically transition to Done when "
     "vulnerabilities are resolved"),
    ("\U0001f6e1\ufe0f  No Duplicates",
     "Tag-based deduplication prevents duplicate work items"),
    ("\U0001f4ca  Security Visibility",
     "Security team sees all vulnerabilities on the ADO board"),
    ("\U0001f3d7\ufe0f  Production-Ready IaC",
     "Bicep templates for repeatable, auditable deployment"),
]

y = 1.5
for title, desc in values:
    add_text_box(slide, 0.8, y, 4, 0.4, title, 18, ORANGE, True)
    add_text_box(slide, 5, y, 7.5, 0.4, desc, 15, LIGHT_GRAY)
    y += 0.7

# ===================== SLIDE 16: Next Steps =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Next Steps", 36, ACCENT_BLUE, True)
add_bullets(slide, 0.8, 1.5, 5.5, 5.5, [
    "1\ufe0f\u20e3  Connect real GitHub repos to the webhook",
    "2\ufe0f\u20e3  Add HMAC signature verification for security",
    "3\ufe0f\u20e3  Configure alert routing to specific teams/areas",
    "4\ufe0f\u20e3  Add Slack/Teams notifications for critical alerts",
    "5\ufe0f\u20e3  Set up monitoring & alerting on the Logic App",
    "6\ufe0f\u20e3  Extend to auto-assign based on code ownership",
], 17, LIGHT_GRAY)

add_text_box(slide, 7, 1.5, 5.5, 0.6, "Resources", 24, ORANGE, True)
add_bullets(slide, 7, 2.5, 5.5, 4, [
    "\U0001f4c2  GitHub Repo:",
    "    github.com/sautalwar/ghas-ado-logic-app",
    "",
    "\u2601\ufe0f  Azure Resource Group:",
    "    rg-ghas-ado-learfield",
    "",
    "\U0001f4cb  ADO Project:",
    "    dev.azure.com/brandsafway1/brandsafway_Engg",
], 14, LIGHT_GRAY)

# ===================== SLIDE 17: Thank You =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 1, 2.0, 11, 1.2,
             "Thank You", 48, ACCENT_BLUE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.5, 11, 0.8,
             "GHAS + Azure Logic Apps + Azure DevOps",
             24, WHITE, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.5, 11, 0.8,
             "Automated security vulnerability tracking \u2014 "
             "from detection to resolution",
             18, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 6.0, 11, 0.5,
             "Questions?", 28, ORANGE, True, PP_ALIGN.CENTER)

# Save
output = "docs/GHAS-ADO-Logic-App-Demo.pptx"
prs.save(output)
print(f"Presentation saved: {output}")
print(f"Total slides: {len(prs.slides)}")
print(f"File size: {os.path.getsize(output) / 1024:.0f} KB")
